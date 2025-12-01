import os
import sys
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from tkinter.ttk import Style
from pptx import Presentation
import re
import platform
from datetime import datetime, timedelta

def to_proper_title_case(text):
    """Convert text to title case: First letter of every word capitalized, rest lowercase."""
    return ' '.join(word.capitalize() for word in text.split())

def convert_to_song_title():
    try:
        # Get the selected text
        highlighted_text = lyrics_text.selection_get()
        if not highlighted_text.strip():
            messagebox.showwarning("No Text Selected", "Please select text to convert.")
            return

        # Remove words in brackets (and the brackets)
        text_no_brackets = re.sub(r'\s*\[.*?\]|\s*\(.*?\)', '', highlighted_text)

        # Remove hyphen and anything to the right (includes different hyphen-like characters)
        text_no_hyphen = re.split(r'[-–—−]', text_no_brackets)[0].strip()

        # Convert to title case
        formatted_text = to_proper_title_case(text_no_hyphen)

        # Format as TITLE (formatted_text)
        formatted_title = f"TITLE ({formatted_text})"

        # Replace the selected text with the formatted title
        lyrics_text.delete(tk.SEL_FIRST, tk.SEL_LAST)
        lyrics_text.insert(tk.INSERT, formatted_title)
    except tk.TclError:
        messagebox.showwarning("No Text Selected", "Please select text to convert.")

# Helper function to get the base directory for bundled or standalone script
def resource_path(relative_path):
    """Get the absolute path to a resource, works for dev and for PyInstaller."""
    base_path = getattr(sys, "_MEIPASS", os.path.dirname(os.path.abspath(__file__)))
    return os.path.join(base_path, relative_path)

# Updated paths for bundled resources
pptx_template_path = resource_path("Central Mega Template v1.pptx")
icon_path = resource_path("icon.png")
logo_path = resource_path("logo.png")

# Function to sanitize file names
def sanitize_file_name(name):
    """Remove invalid characters from the file name."""
    # Remove characters invalid for Windows, macOS, and Linux file systems
    return re.sub(r'[<>:"/\\|?*]', '', name)

# Calculate the next occurrence of a target day (e.g., Sunday or Wednesday).
def get_next_day(start_day, target_day):
    """Calculate the next occurrence of a target day (e.g., Sunday or Wednesday)."""
    days_ahead = target_day - start_day.weekday()
    if days_ahead <= 0:  # Target day is in the next week
        days_ahead += 7
    return start_day + timedelta(days=days_ahead)

# Update the default file name based on the selected service type.
def update_file_name(*args):
    """Update the default file name based on the selected service type."""
    today = datetime.now()
    if service_type_var.get() == "Sunday":
        next_sunday = get_next_day(today, 6)  # Sunday is 6
        default_name = f"Sunday {next_sunday.strftime('%d %m %Y')}"
    else:  # Midweek
        next_wednesday = get_next_day(today, 2)  # Wednesday is 2
        default_name = f"Midweek {next_wednesday.strftime('%d %m %Y')}"
    file_name_entry.delete(0, tk.END)
    file_name_entry.insert(0, default_name)


def save_and_generate_presentation():
    """Generate the slides and save the presentation."""
    global presentation
    file_name = file_name_entry.get()

    if not file_name.strip():
        messagebox.showwarning("Missing File Name", "Please enter a file name.")
        return

    # Get the lyrics from the text box
    song_lyrics = lyrics_text.get("1.0", tk.END).strip()

    # Clean up lyrics by removing excessive blank lines
    cleaned_lyrics = re.sub(r'\n\s*\n\s*\n+', '\n\n', song_lyrics)

    # Check for reserved words in the lyrics
    reserved_words = re.findall(r'(?:^|\s)(?:chorus|verse|bridge|v\d+|[1-9]\.)(?=\s)', cleaned_lyrics, flags=re.IGNORECASE)
    unique_words = "\n".join(set(word.strip() for word in reserved_words))

    if reserved_words:
        result = messagebox.askquestion(
            "Section Labels Found",
            f"You have left the following section label(s) in the lyrics:\n\n{unique_words}\n\n"
            "If these are meant to replace the lines of a verse, chorus, or bridge, you must do so manually.\n"
            "If these are section labels for the singers, it is advised to remove these from the slides.\n"
            "This program cannot differentiate between a verse, chorus, or a bridge.\n\n"
            "Are you sure the lyrics are correct?\n"
            "Would you like to proceed?",
            icon="warning",
            type="yesno"
        )
        if result == "no":
            return

    presentation = Presentation(pptx_template_path)

    # Add the initial slide
    initial_slide_layout = [layout for layout in presentation.slide_layouts if layout.name == "Initial Slide"][0]
    initial_slide = presentation.slides.add_slide(initial_slide_layout)
    for placeholder in initial_slide.placeholders:
        if placeholder.placeholder_format.idx == 10:  # Match the index from the inspection
            if service_type_var.get() == "Sunday":
                placeholder.text = "Our service begins shortly at 11 am"
            else:
                placeholder.text = "Our service begins shortly at 7 pm"

    # Split lyrics into sections based on keywords
    lines = re.split(r'\n\s*\n', cleaned_lyrics)
    sections = []
    current_section = []

    # Case-sensitive keywords
    keywords = ["WELCOME/PRAYER", "COMMUNION", "SERMON", "CLOSE", "CONTRIBUTION", "TITLE"]

    for line in lines:
        if any(line.startswith(keyword + ":") or line.startswith(keyword + " ") for keyword in keywords):
            if current_section:
                sections.append(current_section)
                current_section = []
        current_section.append(line)

    if current_section:
        sections.append(current_section)

    # Process each section
    for section in sections:
        # Check for the first line being a keyword
        match = re.match(r'^(WELCOME/PRAYER|COMMUNION|SERMON|CLOSE|CONTRIBUTION|TITLE)\s*\((.*?)\)$', section[0])
        if match:
            keyword, content = match.groups()
            content = content.strip() if content else ""

            if keyword == "WELCOME/PRAYER":
                welcome_slide_layout = [layout for layout in presentation.slide_layouts if layout.name == "Welcome"][0]
                welcome_slide = presentation.slides.add_slide(welcome_slide_layout)
                for placeholder in welcome_slide.placeholders:
                    if placeholder.placeholder_format.idx == 10:
                        placeholder.text = content

            elif keyword == "COMMUNION":
                communion_slide_layout = [layout for layout in presentation.slide_layouts if layout.name == "Communion"][0]
                communion_slide = presentation.slides.add_slide(communion_slide_layout)
                for placeholder in communion_slide.placeholders:
                    if placeholder.placeholder_format.idx == 10:
                        placeholder.text = content

            elif keyword == "SERMON":
                message_slide_layout = [layout for layout in presentation.slide_layouts if layout.name == "Message"][0]
                message_slide = presentation.slides.add_slide(message_slide_layout)
                for placeholder in message_slide.placeholders:
                    if placeholder.placeholder_format.idx == 10:
                        placeholder.text = content

            elif keyword == "CLOSE":
                close_slide_layout = [layout for layout in presentation.slide_layouts if layout.name == "Close"][0]
                close_slide = presentation.slides.add_slide(close_slide_layout)
                for placeholder in close_slide.placeholders:
                    if placeholder.placeholder_format.idx == 10:
                        placeholder.text = content

            elif keyword == "CONTRIBUTION":
                contribution_slide_layout = [layout for layout in presentation.slide_layouts if layout.name == "Contribution"][0]
                contribution_slide = presentation.slides.add_slide(contribution_slide_layout)
                for placeholder in contribution_slide.placeholders:
                    if placeholder.placeholder_format.idx == 10:
                        placeholder.text = content

                contribution_details_slide_layout = [layout for layout in presentation.slide_layouts if layout.name == "Contribution Details"][0]
                contribution_details_slide = presentation.slides.add_slide(contribution_details_slide_layout)

            elif keyword == "TITLE":
                song_title_slide_layout = [layout for layout in presentation.slide_layouts if layout.name == "Song Title"][0]
                song_title_slide = presentation.slides.add_slide(song_title_slide_layout)
                for placeholder in song_title_slide.placeholders:
                    if placeholder.placeholder_format.idx == 0:
                        placeholder.text = content

        # Process the lyrics in the section
        lyrics_lines = section[1:]  # Exclude the keyword line
        previous_idx10 = None

        for line in lyrics_lines:
            if previous_idx10 is None:
                previous_idx10 = line
                continue

            song_lyrics_slide_layout = [layout for layout in presentation.slide_layouts if layout.name == "Song Lyrics"][0]
            song_lyrics_slide = presentation.slides.add_slide(song_lyrics_slide_layout)

            for placeholder in song_lyrics_slide.placeholders:
                if placeholder.placeholder_format.idx == 10:
                    placeholder.text = previous_idx10
                elif placeholder.placeholder_format.idx == 11:
                    placeholder.text = line

            previous_idx10 = line

        # Handle the last lyric slide in the section
        if previous_idx10:
            final_lyrics_slide_layout = [layout for layout in presentation.slide_layouts if layout.name == "Song Lyrics"][0]
            final_lyrics_slide = presentation.slides.add_slide(final_lyrics_slide_layout)

            for placeholder in final_lyrics_slide.placeholders:
                if placeholder.placeholder_format.idx == 10:
                    placeholder.text = previous_idx10  # Move idx 11 to idx 10
                elif placeholder.placeholder_format.idx == 11:
                    placeholder.text = ""  # Leave idx 11 empty

    # Add the ending slide
    ending_slide_layout = [layout for layout in presentation.slide_layouts if layout.name == "Ending"][0]
    ending_slide = presentation.slides.add_slide(ending_slide_layout)

    # Automatically sanitize the file name
    sanitized_file_name = sanitize_file_name(file_name.strip())  # Clean up whitespace and invalid characters
    default_name = f"{sanitized_file_name}.pptx" if sanitized_file_name else "Presentation.pptx"

    # Ask the user for the file save location
    file_path = filedialog.asksaveasfilename(defaultextension=".pptx",
                                             initialfile=default_name,
                                             filetypes=[("PowerPoint Files", "*.pptx")])
    if file_path:
        presentation.save(file_path)
        messagebox.showinfo("Success", "Presentation saved successfully!")


def handle_undo(event=None):
    """Perform undo operation in the lyrics text widget."""
    lyrics_text.edit_undo()
    return "break"

def handle_redo(event=None):
    """Perform redo operation in the lyrics text widget."""
    lyrics_text.edit_redo()
    return "break"

def show_help():
    """Display help information in a dialog box."""
    help_text = (
        "How to Use Central PPT Slide Maker:\n\n"
        "1. Enter the file name in the 'File Name' field.\n"
        "2. Choose the service type: 'Sunday' or 'Midweek'.\n"
        "3. Paste the song lyrics into the text box labeled 'Order Of Service'.\n"
        "   - Separate slides with new paragraphs.\n"
        "   - Remove any section labels (e.g., 'CHORUS', 'Verse 1').\n\n"
        "4. Use the 'Copy' and 'Paste' buttons to manage text:\n"
        "   - Select text in the 'Order Of Service' text box and click 'Copy' to copy it to the clipboard. Ctrl+C works too.\n"
        "   - Place the cursor where you want to paste text and click 'Paste' to insert text from the clipboard. Ctrl+V works too.\n\n"
        "5. Use the 'Format Selection To Song Title' button to format song titles:\n"
        "   - Select the text in the 'Order Of Service' text box.\n"
        "   - Click 'Format Selection To Song Title'.\n"
        "   - This will format the selection as 'TITLE (Formatted Text)'.\n\n"
        "6. Add title slides using the following keywords (In ALL CAPS):\n"
        "   - WELCOME/PRAYER\n"
        "   - COMMUNION\n"
        "   - SERMON\n"
        "   - CLOSE\n"
        "   - CONTRIBUTION (Includes contribution details slide)\n"
        "   - TITLE (Creates a song title slide)\n\n"
        "   - Each keyword should be followed by the speaker's name in brackets, e.g.:\n"
        "     SERMON (Daniel Marie)\n"
        "     COMMUNION (Jessica Marie).\n"
        "   - If a keyword is not written correctly (typo, not capitalized properly, etc.), the program will not treat it as a title and instead as a song lyric.\n\n"
        "7. Save the presentation:\n"
        "   - Click the 'Save Presentation' button to generate and save the PowerPoint file.\n\n"
        "8. Use 'Undo' and 'Redo' buttons for editing the lyrics text box.\n"
        "   - You can also use shortcuts:\n"
        "     - Ctrl+Z for Undo.\n"
        "     - Ctrl+Y or Ctrl+Shift+Z for Redo.\n\n"
        "Central PPT Slide Maker 2.0.1"
    )
    messagebox.showinfo("Help - Central PPT Slide Maker", help_text)




# Create the GUI window
root = tk.Tk()
root.title("Central PPT Slide Maker")
root.configure(bg="#f7f1e3")
root.option_add("*Font", "Helvetica 11")  # Set global font

# Configure ttk style for Radiobutton to match the background
style = Style()
style.configure("Custom.TRadiobutton", background="#f7f1e3", font=("Helvetica", 11))

# Change program icon
try:
    root.iconbitmap(icon_path)
except Exception as e:
    print(f"Error loading icon: {e}")

# Add a menu bar with a Help button
menu_bar = tk.Menu(root)
help_menu = tk.Menu(menu_bar, tearoff=0)
help_menu.add_command(label="How to Use", command=show_help)
menu_bar.add_cascade(label="Help", menu=help_menu)
root.config(menu=menu_bar)

# Add a logo at the top
try:
    logo = tk.PhotoImage(file=logo_path)
    logo_label = tk.Label(root, image=logo, bg="#f7f1e3")
except Exception as e:
    print(f"Error loading logo: {e}")
    logo_label = None # Ensure logo_label exists even if loading fails

# File name input row
file_name_frame = tk.Frame(root, bg="#f7f1e3")
tk.Label(file_name_frame, text="File Name:", bg="#f7f1e3").pack(side=tk.LEFT, padx=5)
file_name_entry = ttk.Entry(file_name_frame, width=40)
file_name_entry.pack(side=tk.LEFT, padx=5)

# Service type toggle row
service_type_frame = tk.Frame(root, bg="#f7f1e3")
tk.Label(service_type_frame, text="Service Type:", bg="#f7f1e3").pack(side=tk.LEFT, padx=5)
service_type_var = tk.StringVar(value="Sunday")
sunday_radio = ttk.Radiobutton(service_type_frame, text="Sunday", variable=service_type_var, value="Sunday", style="Custom.TRadiobutton")
midweek_radio = ttk.Radiobutton(service_type_frame, text="Midweek", variable=service_type_var, value="Midweek", style="Custom.TRadiobutton")
sunday_radio.pack(side=tk.LEFT, padx=5)
midweek_radio.pack(side=tk.LEFT, padx=5)
service_type_var.trace_add("write", lambda *args: update_file_name())
update_file_name()

# Lyrics input
lyrics_label = tk.Label(root, text="Order Of Service:", bg="#f7f1e3")
lyrics_frame = tk.Frame(root, bg="#f7f1e3")
lyrics_text = tk.Text(lyrics_frame, wrap=tk.WORD, undo=True, font=("Helvetica", 12))
lyrics_scrollbar = tk.Scrollbar(lyrics_frame, command=lyrics_text.yview)
lyrics_text.config(yscrollcommand=lyrics_scrollbar.set)
lyrics_text.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=(10, 0))
lyrics_scrollbar.pack(side=tk.RIGHT, fill=tk.Y, padx=(0, 10))

# Button row
button_frame = tk.Frame(root, bg="#f7f1e3")
undo_btn = ttk.Button(button_frame, text="Undo", command=handle_undo)
undo_btn.pack(side=tk.LEFT, padx=(5, 1))
redo_btn = ttk.Button(button_frame, text="Redo", command=handle_redo)
redo_btn.pack(side=tk.LEFT, padx=(1, 10))
copy_btn = ttk.Button(button_frame, text="Copy", command=lambda: root.clipboard_append(lyrics_text.selection_get()))
copy_btn.pack(side=tk.LEFT, padx=(10, 1))
paste_btn = ttk.Button(button_frame, text="Paste", command=lambda: lyrics_text.insert(tk.INSERT, root.clipboard_get()))
paste_btn.pack(side=tk.LEFT, padx=(1, 10))
format_btn = ttk.Button(button_frame, text="Format Selection To Song Title", command=convert_to_song_title)
format_btn.pack(side=tk.LEFT, padx=(10, 5))
save_and_generate_btn = ttk.Button(button_frame, text="Save Presentation", command=save_and_generate_presentation)
save_and_generate_btn.pack(side=tk.RIGHT, padx=5)


# Pack the button frame and anchor it to the BOTTOM of the window first.
button_frame.pack(side=tk.BOTTOM, pady=10, fill=tk.X)

# Pack all the top elements in order from the top down.
if logo_label:
    logo_label.pack(side=tk.TOP, pady=10)
file_name_frame.pack(side=tk.TOP, pady=10, anchor="center")
service_type_frame.pack(side=tk.TOP, pady=10, anchor="center")
lyrics_label.pack(side=tk.TOP, anchor="center", pady=(5, 0))

# Finally, pack the lyrics frame.
lyrics_frame.pack(side=tk.TOP, fill=tk.BOTH, expand=True, pady=(2, 10))


# Bind undo/redo shortcuts
if platform.system() == "Darwin":  # macOS
    root.bind("<Command-z>", handle_undo)  # Cmd+Z for undo
    root.bind("<Command-Shift-z>", handle_redo)  # Cmd+Shift+Z for redo
else:  # Windows/Linux
    root.bind("<Control-z>", handle_undo)  # Ctrl+Z for undo
    root.bind("<Control-y>", handle_redo)  # Ctrl+Y for redo
    root.bind("<Control-Shift-Z>", handle_redo)  # Ctrl+Shift+Z for redo (capital Z is interpreted differently)

root.mainloop()
